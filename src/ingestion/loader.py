

from pathlib import Path
from typing import List, Optional, Dict, Any

import aiofiles

from ..utils.exceptions import IngestionError


class Document:

    def __init__(
        self,
        content: str,
        metadata: Optional[Dict[str, Any]] = None,
        doc_id: Optional[str] = None,
    ):
        self.content = content
        self.metadata = metadata or {}
        self.doc_id = doc_id or self.metadata.get("source", "unknown")

    def __repr__(self) -> str:
        return f"Document(id={self.doc_id}, length={len(self.content)})"


class DocumentLoader:

    async def load(self, source: str) -> List[Document]:
        raise NotImplementedError


class TextLoader(DocumentLoader):

    async def load(self, file_path: str) -> List[Document]:
        try:
            async with aiofiles.open(
                file_path,
                "r",
                encoding="utf-8",
            ) as file:
                content = await file.read()

            return [
                Document(
                    content=content,
                    metadata={
                        "source": file_path,
                        "type": "text",
                    },
                )
            ]

        except Exception as error:
            raise IngestionError(
                f"Failed to load text file: {file_path}",
                original_error=error,
            )
class PDFLoader(DocumentLoader):
    

    async def load(self, file_path: str) -> List[Document]:
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            documents = []

            for page_number, page in enumerate(reader.pages):
                content = page.extract_text() or ""

                if content.strip():
                    documents.append(
                        Document(
                            content=content,
                            metadata={
                                "source": file_path,
                                "type": "pdf",
                                "page": page_number + 1,
                            },
                        )
                    )

            return documents

        except Exception as error:
            raise IngestionError(
                f"Failed to load PDF file: {file_path}",
                original_error=error,
            )

class MarkdownLoader(DocumentLoader):

    async def load(self, file_path: str) -> List[Document]:
        try:
            async with aiofiles.open(
                file_path,
                "r",
                encoding="utf-8",
            ) as file:
                content = await file.read()

            return [
                Document(
                    content=content,
                    metadata={
                        "source": file_path,
                        "type": "markdown",
                    },
                )
            ]

        except Exception as error:
            raise IngestionError(
                f"Failed to load Markdown file: {file_path}",
                original_error=error,
            )


class HTMLLoader(DocumentLoader):

    async def load(self, file_path: str) -> List[Document]:
        try:
            from bs4 import BeautifulSoup

            async with aiofiles.open(
                file_path,
                "r",
                encoding="utf-8",
            ) as file:
                content = await file.read()

            soup = BeautifulSoup(content, "html.parser")

            return [
                Document(
                    content=soup.get_text(separator="\n", strip=True),
                    metadata={
                        "source": file_path,
                        "type": "html",
                    },
                )
            ]

        except Exception as error:
            raise IngestionError(
                f"Failed to load HTML file: {file_path}",
                original_error=error,
            )
class CSVLoader(DocumentLoader):
    

    async def load(self, file_path: str) -> List[Document]:
        try:
            import csv

            async with aiofiles.open(
                file_path,
                "r",
                encoding="utf-8",
            ) as file:
                content = await file.read()

            rows = list(csv.reader(content.splitlines()))

            if not rows:
                return []

            headers = rows[0]
            documents = []

            for index, row in enumerate(rows[1:], start=1):
                row_data = dict(zip(headers, row))

                documents.append(
                    Document(
                        content=str(row_data),
                        metadata={
                            "source": file_path,
                            "type": "csv",
                            "row": index,
                        },
                    )
                )

            return documents

        except Exception as error:
            raise IngestionError(
                f"Failed to load CSV file: {file_path}",
                original_error=error,
            )


class DOCXLoader(DocumentLoader):

    async def load(self, file_path: str) -> List[Document]:
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)

            content = "\n".join(
                paragraph.text
                for paragraph in doc.paragraphs
                if paragraph.text.strip()
            )

            return [
                Document(
                    content=content,
                    metadata={
                        "source": file_path,
                        "type": "docx",
                    },
                )
            ]

        except Exception as error:
            raise IngestionError(
                f"Failed to load DOCX file: {file_path}",
                original_error=error,
            )

class ExcelLoader(DocumentLoader):

    async def load(self, file_path: str) -> List[Document]:
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(
                file_path,
                read_only=True,
                data_only=True,
            )

            documents = []

            for sheet in workbook.worksheets:
                rows = sheet.iter_rows(values_only=True)

                for row_number, row in enumerate(rows, start=1):
                    values = [
                        str(value)
                        for value in row
                        if value is not None
                    ]

                    if not values:
                        continue

                    documents.append(
                        Document(
                            content=" | ".join(values),
                            metadata={
                                "source": file_path,
                                "type": "xlsx",
                                "sheet": sheet.title,
                                "row": row_number,
                            },
                        )
                    )

            workbook.close()
            return documents

        except Exception as error:
            raise IngestionError(
                f"Failed to load Excel file: {file_path}",
                original_error=error,
            )

class PPTXLoader(DocumentLoader):

    async def load(self, file_path: str) -> List[Document]:
        try:
            from pptx import Presentation

            presentation = Presentation(file_path)
            documents = []

            for slide_number, slide in enumerate(
                presentation.slides,
                start=1,
            ):
                texts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())

                content = "\n".join(texts)

                if content:
                    documents.append(
                        Document(
                            content=content,
                            metadata={
                                "source": file_path,
                                "type": "pptx",
                                "slide": slide_number,
                            },
                        )
                    )

            return documents

        except Exception as error:
            raise IngestionError(
                f"Failed to load PowerPoint file: {file_path}",
                original_error=error,
            )

class JSONLoader(DocumentLoader):

    async def load(self, file_path: str) -> List[Document]:
        try:
            import json

            async with aiofiles.open(
                file_path,
                "r",
                encoding="utf-8",
            ) as file:
                content = await file.read()

            data = json.loads(content)

            return [
                Document(
                    content=json.dumps(data, indent=2),
                    metadata={
                        "source": file_path,
                        "type": "json",
                    },
                )
            ]

        except Exception as error:
            raise IngestionError(
                f"Failed to load JSON file: {file_path}",
                original_error=error,
            )

class ImageLoader(DocumentLoader):

    async def load(self, file_path: str) -> List[Document]:
        try:
            from PIL import Image
            import pytesseract

            image = Image.open(file_path)
            content = pytesseract.image_to_string(image)

            return [
                Document(
                    content=content,
                    metadata={
                        "source": file_path,
                        "type": "image",
                    },
                )
            ]

        except Exception as error:
            raise IngestionError(
                f"Failed to load image file: {file_path}",
                original_error=error,
            )

def get_loader(file_path: str) -> DocumentLoader:

    suffix = Path(file_path).suffix.lower()

    loaders = {
        ".txt": TextLoader(),
        ".pdf": PDFLoader(),
        ".md": MarkdownLoader(),
        ".markdown": MarkdownLoader(),
        ".html": HTMLLoader(),
        ".htm": HTMLLoader(),
        ".csv": CSVLoader(),
        ".docx": DOCXLoader(),
        ".xlsx": ExcelLoader(),
        ".pptx": PPTXLoader(),
        ".json": JSONLoader(),
        ".png": ImageLoader(),
        ".jpg": ImageLoader(),
        ".jpeg": ImageLoader(),
    }

    if suffix not in loaders:
        raise IngestionError(
            f"Unsupported file type: {suffix}"
        )

    return loaders[suffix]


async def load_document(file_path: str) -> List[Document]:
    loader = get_loader(file_path)
    return await loader.load(file_path)